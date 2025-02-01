import os
import fitz  # PyMuPDF for PDF processing
from fastapi import FastAPI, UploadFile, HTTPException
from dotenv import load_dotenv
from diffusers import StableDiffusionXLPipeline
from elevenlabs import generate
import requests
from pptx import Presentation
import uvicorn
from huggingface_hub import login
import ollama
import logging
from fpdf import FPDF
import torch
from datetime import datetime

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[logging.FileHandler("app.log"), logging.StreamHandler()]
)

# Load environment variables
load_dotenv()

# Create a directory for temporary files
TEMP_DIR = "temp_files"
os.makedirs(TEMP_DIR, exist_ok=True)

# Get API keys securely
ELEVENLABS_API_KEY = os.getenv("ELEVENLABS_API_KEY")
HUGGINGFACE_API_TOKEN = os.getenv("HUGGINGFACE_API_TOKEN")

# Ensure API keys are set
if not ELEVENLABS_API_KEY or not HUGGINGFACE_API_TOKEN:
    logging.error("API keys are not set correctly.")
    raise ValueError("API keys are missing.")

# Authenticate with Hugging Face Hub
login(token=HUGGINGFACE_API_TOKEN)

class SummaryPDF(FPDF):
    def header(self):
        self.set_font('Arial', 'B', 15)
        self.cell(0, 10, 'Research Paper Summary', 0, 1, 'C')
        self.ln(10)
        
    def footer(self):
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.cell(0, 10, f'Generated on {datetime.now().strftime("%Y-%m-%d %H:%M")} - Page {self.page_no()}', 0, 0, 'C')

app = FastAPI()

def initialize_stable_diffusion():
    """Initialize Stable Diffusion with optimal settings"""
    try:
        # Check if CUDA is available
        if not torch.cuda.is_available():
            raise RuntimeError("CUDA is not available. Please ensure you have a CUDA-capable GPU.")
        
        # Initialize the pipeline with proper settings
        pipe = StableDiffusionXLPipeline.from_pretrained(
            "stabilityai/stable-diffusion-xl-base-1.0",
            torch_dtype=torch.float16,  # Ensure float16 for memory efficiency
            use_safetensors=True,
            variant="fp16"
        )
        
        # Move to CUDA device
        pipe = pipe.to("cuda")
        
        # Optional memory optimizations based on available GPU memory
        if torch.cuda.get_device_properties(0).total_memory >= 8 * (1024 ** 3):  # 8GB or more
            try:
                pipe.enable_xformers_memory_efficient_attention()
            except Exception as e:
                logging.warning(f"Could not enable xformers: {e}")
                # Fallback to standard attention mechanism
                pipe.enable_attention_slicing()
        else:
            # For GPUs with less memory, use these optimizations
            pipe.enable_attention_slicing()
            pipe.enable_sequential_cpu_offload()
        
        return pipe
    except Exception as e:
        logging.error(f"Error initializing Stable Diffusion: {str(e)}")
        # Provide more specific error messages based on common issues
        if "CUDA" in str(e):
            raise HTTPException(
                status_code=500,
                detail="GPU initialization failed. Please ensure CUDA is properly installed and a compatible GPU is available."
            )
        elif "memory" in str(e).lower():
            raise HTTPException(
                status_code=500,
                detail="Insufficient GPU memory. Try reducing batch size or image dimensions."
            )
        else:
            raise HTTPException(status_code=500, detail=f"Failed to initialize Stable Diffusion: {str(e)}")

def generate_graphical_abstract(summary, pipe):
    try:
        # Create a more focused prompt for research visualization
        prompt = f"""Create a scientific graphical abstract visualization:
        A clean, professional diagram showing:
        {summary[:300]}
        Style: Modern scientific illustration, minimalist, clear layout, professional colors
        Include: Relevant scientific symbols, data visualization elements, and clear visual hierarchy
        """
        
        negative_prompt = "text, words, blurry, low quality, distorted, messy, cluttered"
        
        # Generate image with improved parameters and error handling
        try:
            output = pipe(
                prompt=prompt,
                negative_prompt=negative_prompt,
                num_inference_steps=30,  # Reduced steps for better performance
                guidance_scale=7.5,
                height=512,  # Reduced size for memory efficiency
                width=512,
                generator=torch.manual_seed(42)  # For reproducibility
            )
        except torch.cuda.OutOfMemoryError:
            # Fallback to smaller image size if OOM error occurs
            output = pipe(
                prompt=prompt,
                negative_prompt=negative_prompt,
                num_inference_steps=20,
                guidance_scale=7.5,
                height=384,
                width=384,
                generator=torch.manual_seed(42)
            )
        
        if not output.images:
            raise ValueError("No images generated")
            
        image = output.images[0]
        graphical_abstract_path = os.path.join(TEMP_DIR, "graphical_abstract.png")
        image.save(graphical_abstract_path, quality=95)
        
        return graphical_abstract_path
    except Exception as e:
        logging.error(f"Graphical Abstract Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate graphical abstract: {str(e)}")

@app.get("/")
def read_root():
    return {"message": "FastAPI server is running!"}

# Extract text from PDF
def extract_text_from_pdf(pdf_path):
    try:
        with fitz.open(pdf_path) as doc:
            return "".join(page.get_text("text") for page in doc)
    except Exception as e:
        logging.error(f"PDF Extraction Error: {str(e)}")
        raise HTTPException(status_code=400, detail=f"PDF Extraction Error: {str(e)}")

# Ollama Summary
def ollama_summary(text, summary_length="medium"):
    try:
        response = ollama.generate(
            model="llama2",
            prompt=f"Summarize the following research paper in {summary_length} length, providing key findings, methodology, conclusions, and implications: {text}",
            options={"max_tokens": 1000}
        )
        logging.info("Ollama summary generated successfully.")
        return response["response"] if "response" in response else ""
    except Exception as e:
        logging.error(f"Ollama Summary Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Ollama Summary Error: {str(e)}")

def format_summary_sections(summary):
    """Format the summary into structured sections with improved parsing"""
    sections = {
        "Key Findings": [],
        "Methodology": [],
        "Conclusions": [],
        "Implications": []
    }
    
    try:
        # Modify the prompt to get more structured output
        categorized = ollama.generate(
            model="llama2",
            prompt=f"""Please analyze this summary and provide a detailed breakdown in the following format:

KEY FINDINGS:
- First key finding
- Second key finding
(etc.)

METHODOLOGY:
- First methodological point
- Second methodological point
(etc.)

CONCLUSIONS:
- First conclusion
- Second conclusion
(etc.)

IMPLICATIONS:
- First implication
- Second implication
(etc.)

Summary to analyze: {summary}""",
            options={"max_tokens": 2000}
        )
        
        # Improved parsing logic
        current_section = None
        response_lines = categorized["response"].split('\n')
        
        for line in response_lines:
            line = line.strip()
            if not line:
                continue
                
            # Check for section headers
            upper_line = line.upper()
            if "KEY FINDINGS:" in upper_line:
                current_section = "Key Findings"
            elif "METHODOLOGY:" in upper_line:
                current_section = "Methodology"
            elif "CONCLUSIONS:" in upper_line:
                current_section = "Conclusions"
            elif "IMPLICATIONS:" in upper_line:
                current_section = "Implications"
            # Add content if we're in a section and the line starts with a bullet point or dash
            elif current_section and (line.startswith('-') or line.startswith('•')):
                # Remove the bullet point/dash and clean up the text
                content = line.lstrip('- •').strip()
                if content:  # Only add non-empty content
                    sections[current_section].append(content)
                    
        # Ensure each section has at least some content
        for section in sections:
            if not sections[section]:
                # Generate specific content for empty sections
                section_content = ollama.generate(
                    model="llama2",
                    prompt=f"Based on this summary, provide 2-3 points specifically for the {section} section: {summary}",
                    options={"max_tokens": 500}
                )
                # Parse the response and add points
                for line in section_content["response"].split('\n'):
                    line = line.strip()
                    if line and not line.upper().endswith(':'):
                        sections[section].append(line.lstrip('- •'))
    
    except Exception as e:
        logging.error(f"Error in format_summary_sections: {str(e)}")
        # Fallback: generate content for each section independently
        for section in sections:
            try:
                section_content = ollama.generate(
                    model="llama2",
                    prompt=f"Provide 2-3 key points for the {section} section based on this summary: {summary}",
                    options={"max_tokens": 500}
                )
                # Parse the response and add points
                for line in section_content["response"].split('\n'):
                    line = line.strip()
                    if line and not line.upper().endswith(':'):
                        sections[section].append(line.lstrip('- •'))
            except:
                sections[section].append("Content generation failed for this section")
    
    return sections

class SummaryPDF(FPDF):
    def header(self):
        self.set_font('Arial', 'B', 15)
        self.cell(0, 10, 'Research Paper Summary', 0, 1, 'C')
        self.ln(10)
        
    def footer(self):
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.cell(0, 10, f'Generated on {datetime.now().strftime("%Y-%m-%d %H:%M")} - Page {self.page_no()}', 0, 0, 'C')

def save_summary_to_pdf(summary, output_path):
    try:
        pdf = SummaryPDF()
        pdf.add_page()
        
        # Format the summary into sections
        sections = format_summary_sections(summary)
        
        # Add formatted content using standard dash
        for section, points in sections.items():
            # Section header
            pdf.set_font('Arial', 'B', 14)
            pdf.cell(0, 10, section, ln=True)
            pdf.ln(5)
            
            # Section content
            pdf.set_font('Arial', '', 12)
            for point in points:
                # Clean the text of any problematic characters
                clean_point = point.encode('ascii', 'replace').decode()
                # Add content with dash instead of bullet
                pdf.multi_cell(0, 10, f"- {clean_point}")
            pdf.ln(5)
        
        pdf.output(output_path)
        logging.info(f"Summary saved to {output_path}")
        return output_path
    except Exception as e:
        logging.error(f"Error saving summary to PDF: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error saving summary to PDF: {str(e)}")
def generate_graphical_abstract(summary, pipe):
    try:
        # Create a more focused prompt for research visualization
        prompt = f"""Create a scientific graphical abstract visualization:
        A clean, professional diagram showing:
        {summary[:300]}
        Style: Modern scientific illustration, minimalist, clear layout, professional colors
        Include: Relevant scientific symbols, data visualization elements, and clear visual hierarchy
        """
        
        negative_prompt = "text, words, blurry, low quality, distorted, messy, cluttered"
        
        # Generate image with improved parameters
        output = pipe(
            prompt=prompt,
            negative_prompt=negative_prompt,
            num_inference_steps=50,
            guidance_scale=7.5,
            height=768,
            width=768
        )
        
        if not output.images:
            raise ValueError("No images generated")
            
        image = output.images[0]
        graphical_abstract_path = os.path.join(TEMP_DIR, "graphical_abstract.png")
        image.save(graphical_abstract_path, quality=95)
        
        return graphical_abstract_path
    except Exception as e:
        logging.error(f"Graphical Abstract Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate graphical abstract: {str(e)}")

# AI Voiceover
def generate_voice(summary):
    try:
        audio = generate(
            text=summary,
            voice="Lily",
            model="eleven_monolingual_v1",
            api_key=ELEVENLABS_API_KEY
        )
        voiceover_path = os.path.join(TEMP_DIR, "voiceover.mp3")
        with open(voiceover_path, "wb") as f:
            f.write(audio)
        return voiceover_path
    except Exception as e:
        logging.error(f"Voiceover Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Voiceover Error: {str(e)}")

# Generate Presentation
def generate_presentation(summary):
    try:
        prs = Presentation()
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Research Paper Summary"
        title_slide.placeholders[1].text = "Generated Summary Presentation"
        
        # Summary sections
        sections = format_summary_sections(summary)
        for section, points in sections.items():
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = section
            text_frame = slide.placeholders[1].text_frame
            
            for point in points:
                p = text_frame.add_paragraph()
                p.text = f"• {point}"
                p.level = 0
        
        presentation_path = os.path.join(TEMP_DIR, "presentation.pptx")
        prs.save(presentation_path)
        return presentation_path
    except Exception as e:
        logging.error(f"Presentation Generation Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Presentation Generation Error: {str(e)}")

# FastAPI Route
@app.post("/process-paper/")
async def process_paper(file: UploadFile, summary_length: str = "medium"):
    try:
        # Initialize Stable Diffusion
        pipe = initialize_stable_diffusion()
        
        file_path = os.path.join(TEMP_DIR, "temp_paper.pdf")
        with open(file_path, "wb") as f:
            while chunk := file.file.read(1024 * 1024):
                f.write(chunk)

        text = extract_text_from_pdf(file_path)
        summary = ollama_summary(text, summary_length)
        
        # Generate all outputs
        summary_pdf_path = os.path.join(TEMP_DIR, "summary.pdf")
        save_summary_to_pdf(summary, summary_pdf_path)
        
        graphical_abstract_path = generate_graphical_abstract(summary, pipe)
        voiceover_path = generate_voice(summary)
        presentation_path = generate_presentation(summary)

        return {
            "summary": summary,
            "summary_pdf": summary_pdf_path,
            "graphical_abstract": graphical_abstract_path,
            "voiceover": voiceover_path,
            "presentation": presentation_path,
        }
    except Exception as e:
        logging.error(f"Error processing paper: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

# Run FastAPI Server
if __name__ == "__main__":
    uvicorn.run("main:app", host="0.0.0.0", port=8080, reload=True)
